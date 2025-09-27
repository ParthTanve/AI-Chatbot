import speech_recognition as sr
from gtts import gTTS
import tempfile
import os
import webbrowser
from PIL import Image
import fitz
import time
import pandas as pd
import openpyxl
import pptx
import docx
from langchain.memory import StreamlitChatMessageHistory
import google.generativeai as genai
import streamlit as st


def speak_text(text):
    """Generates and plays audio from text using gTTS."""
    try:
        
        tts = gTTS(text=text, lang="en")
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as fp:
            tts.save(fp.name)
            audio_path = fp.name
        
        
        st.audio(audio_path, format="audio/mp3")
        os.remove(audio_path) 
    except Exception as e:
    
        pass 


def speech_to_text_from_mic(audio_bytes):
    """
    Performs speech-to-text on raw audio bytes captured from the browser.
    
    Args:
        audio_bytes (bytes): Raw audio data in WAV format.
    
    Returns:
        str: Recognized text or an error message.
    """
    
    if not audio_bytes or not isinstance(audio_bytes, bytes):
        return "Speech-to-Text Error: No valid audio data received."
    
    try:
        recognizer = sr.Recognizer()
        
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as fp:
            fp.write(audio_bytes)
            audio_path = fp.name
        
        
        
        with sr.AudioFile(audio_path) as source:
            audio = recognizer.record(source)
            text = recognizer.recognize_google(audio)
            
        os.remove(audio_path) 
        return text
    
    except sr.UnknownValueError:
        return "Speech-to-Text Error: Could not understand audio (UnknownValueError)."
    except Exception as e:
        return f"Speech-to-Text Error: {e}"


def extract_text_from_file(file):
    """Extracts text content from various file types."""
    try:
        file_extension = os.path.splitext(file.name)[1].lower()
        
        if file_extension == ".pdf":
            doc = fitz.open(stream=file.read(), filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            return text
        elif file_extension in [".txt", ".py", ".c", ".cpp", ".java", ".js", ".html", ".css"]:
            return file.read().decode("utf-8")
        elif file_extension in [".xlsx", ".xls"]:
            df = pd.read_excel(file)
            return df.to_string()
        elif file_extension in [".pptx", ".ppt"]:
            presentation = pptx.Presentation(file)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        text += shape.text_frame.text
            return text
        elif file_extension == ".docx":
            doc = docx.Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        else:
            return "Unsupported file type."
    except Exception as e:
        return f"File processing error: {e}"


def handle_command(user_input):
    """Handles simple hardcoded commands like 'open website'."""
    commands = {
        "open google": "https://www.google.com",
        "open github": "https://www.github.com"
    }

    user_input = user_input.lower()

    if "open" in user_input:
        for key, url in commands.items():
            if key in user_input:
                webbrowser.open(url)
                return f"Opening {key.replace('open ', '')}..."

        if "notepad" in user_input:
           
            os.startfile("notepad.exe") 
            return "Opening Notepad..."

        words = user_input.split()
        for word in words:
            if "." in word:
                if not word.startswith("http"):
                    word = "https://" + word
                webbrowser.open(word)
                return f"Opening {word}..."

    return None


def ask_gemini(prompt, model, image=None, chat_history=None, retries=3):
    """Handles interaction with the Gemini API, including chat history and image input."""
    full_prompt = ""
    
    if chat_history:
        for msg in chat_history:
            if hasattr(msg, 'type') and hasattr(msg, 'content'):
                if msg.type == "human":
                    full_prompt += f"User: {msg.content}\n"
                elif msg.type == "ai":
                    full_prompt += f"Assistant: {msg.content}\n"
    full_prompt += f"User: {prompt}\nAssistant:"

    for i in range(retries):
        try:
            content_list = [full_prompt]
            if image:
                content_list.append(image)
                
            response = model.generate_content(content_list)

            if response.candidates:
                candidate = response.candidates[0]
                if hasattr(candidate, "content") and candidate.content.parts:
                    texts = [part.text for part in candidate.content.parts if hasattr(part, "text")]
                    if texts:
                        return " ".join(texts)
            return "Sorry, I couldn’t generate a response."
        except Exception as e:
            if i < retries - 1:
                time.sleep(2**(i+1))
            else:
                return f"API Error: {e}"
    return "Failed to get a response after multiple retries."
